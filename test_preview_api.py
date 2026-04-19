import os,sys,tempfile
from io import BytesIO
sys.path.insert(0,r'D:\OperatingProject\project')
from config import Config
td = tempfile.mkdtemp()
Config.DATABASE_PATH = os.path.join(td,'api.db')
Config.STORAGE_PATH = os.path.join(td,'storage')
Config.ENCRYPTION_KEY_FILE = os.path.join(td,'key')
from main import create_app
from app.protection.encryption import EncryptionService
from pptx import Presentation

EncryptionService.reset()
app = create_app(testing=True)
c = app.test_client()
r = c.post('/api/auth/register',json={'username':'api_test','password':'Str0ng!Pass'})
l = c.post('/api/auth/login',json={'username':'api_test','password':'Str0ng!Pass'})
t = l.get_json()['access_token']
prs = Presentation()
s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = 'API Preview Test'
s.placeholders[1].text = 'Checking PDF output'
buf = BytesIO()
prs.save(buf)
u = c.post('/api/files/upload',headers={'Authorization':f'Bearer {t}'},data={'file':(BytesIO(buf.getvalue()),'api_test.pptx')},content_type='multipart/form-data')
fid = u.get_json()['file']['id']
p = c.get(f'/api/files/{fid}/preview?token={t}')
print(f'PREVIEW_STATUS={p.status_code}')
print(f'PREVIEW_CONTENT_TYPE={p.content_type}')
print(f'PREVIEW_IS_PDF={p.content_type=="application/pdf"}')
print(f'PREVIEW_BYTES={len(p.data)}')
