import os
import sys
import tempfile
import unittest
from io import BytesIO
from unittest.mock import patch
import json

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))

from config import Config

_test_dir = tempfile.mkdtemp()
Config.DATABASE_PATH = os.path.join(_test_dir, "test_features.db")
Config.STORAGE_PATH = os.path.join(_test_dir, "storage")
Config.ENCRYPTION_KEY_FILE = os.path.join(_test_dir, "test_features.key")

from main import create_app
from app.protection.encryption import EncryptionService
from app.files.intelligence import FileIntelligenceService


class TestFileFeatureRoutes(unittest.TestCase):

    @classmethod
    def setUpClass(cls):
        EncryptionService.reset()
        cls.app = create_app(testing=True)
        cls.client = cls.app.test_client()

        register_resp = cls.client.post(
            '/api/auth/register',
            json={"username": "feature_user", "password": "Str0ng!Pass"}
        )
        assert register_resp.status_code == 201

        login_resp = cls.client.post(
            '/api/auth/login',
            json={"username": "feature_user", "password": "Str0ng!Pass"}
        )
        assert login_resp.status_code == 200
        cls.token = login_resp.get_json()["access_token"]

        upload_resp = cls.client.post(
            '/api/files/upload',
            headers={"Authorization": f"Bearer {cls.token}"},
            data={"file": (BytesIO(b"Confidential planning notes."), "feature.txt")},
            content_type='multipart/form-data'
        )
        assert upload_resp.status_code == 201
        cls.file_id = upload_resp.get_json()["file"]["id"]

    def test_insights_route(self):
        response = self.client.get(
            f'/api/files/{self.file_id}/insights',
            headers={"Authorization": f"Bearer {self.token}"}
        )
        self.assertEqual(response.status_code, 200)
        body = response.get_json()
        self.assertIn("insights", body)
        self.assertIn("summary", body["insights"])

    def test_groq_ai_insights_route_uses_api_response(self):
        groq_response = {
            "choices": [
                {
                    "message": {
                        "content": json.dumps({
                            "summary": "Groq generated summary for the uploaded document.",
                            "keywords": ["groq", "summary", "document"],
                            "tags": ["ai", "cloud"],
                            "sensitivity": "low",
                            "suggested_actions": ["Keep encrypted at rest", "Share only with least-privilege access"],
                        })
                    }
                }
            ]
        }

        class FakeResponse:
            def __enter__(self):
                return self

            def __exit__(self, exc_type, exc, tb):
                return False

            def read(self):
                return json.dumps(groq_response).encode('utf-8')

        with patch.object(FileIntelligenceService, 'build_insights', wraps=FileIntelligenceService.build_insights):
            with patch('app.files.intelligence.Config.GROQ_API_KEY', 'dummy-key'):
                with patch('app.files.intelligence.GROQ_SDK_AVAILABLE', False):
                    with patch('app.files.intelligence.urlopen', return_value=FakeResponse()):
                        response = self.client.get(
                            f'/api/files/{self.file_id}/insights',
                            headers={"Authorization": f"Bearer {self.token}"}
                        )

        self.assertEqual(response.status_code, 200)
        body = response.get_json()
        self.assertEqual(body["insights"]["source"], "groq")
        self.assertEqual(body["insights"]["engine"], "groq:llama-3.3-70b-versatile")
        self.assertIn("Groq generated summary", body["insights"]["summary"])

    def test_legacy_ppt_insights_route_returns_200(self):
        ppt_like_payload = (
            b"PowerPoint Document\x00\x01"
            b"Operating System Synchronization\x00"
            b"Process\x00Thread\x00Semaphore\x00Monitor\x00"
        )

        upload_resp = self.client.post(
            '/api/files/upload',
            headers={"Authorization": f"Bearer {self.token}"},
            data={"file": (BytesIO(ppt_like_payload), "legacy.ppt")},
            content_type='multipart/form-data'
        )
        self.assertEqual(upload_resp.status_code, 201)
        file_id = upload_resp.get_json()["file"]["id"]

        response = self.client.get(
            f'/api/files/{file_id}/insights',
            headers={"Authorization": f"Bearer {self.token}"}
        )
        self.assertEqual(response.status_code, 200)
        body = response.get_json()
        self.assertIn("insights", body)
        self.assertIn("summary", body["insights"])
        self.assertIn(body["insights"]["source"], ["groq", "groq_unavailable"])

    def test_legacy_ppt_preview_route_returns_html(self):
        ppt_like_payload = (
            b"PowerPoint Document\x00\x01"
            b"OS Process Synchronization\x00"
            b"Critical Section\x00Semaphore\x00Monitor\x00"
        )

        upload_resp = self.client.post(
            '/api/files/upload',
            headers={"Authorization": f"Bearer {self.token}"},
            data={"file": (BytesIO(ppt_like_payload), "lecture.ppt")},
            content_type='multipart/form-data'
        )
        self.assertEqual(upload_resp.status_code, 201)
        file_id = upload_resp.get_json()["file"]["id"]

        preview_resp = self.client.get(
            f'/api/files/{file_id}/preview?token={self.token}'
        )
        self.assertEqual(preview_resp.status_code, 200)
        if 'application/json' in preview_resp.content_type:
            body = preview_resp.get_json()
            self.assertIn("preview", body)
            self.assertIn("html", body["preview"])
            self.assertIn("Legacy PowerPoint Preview", body["preview"]["html"])
        else:
            self.assertIn('application/pdf', preview_resp.content_type)

    def test_legacy_ppt_binary_garbage_is_filtered(self):
        noisy_payload = (
            b"IDATx^\x00\x01\x02\x03"
            b"p9dy9\x00\x00QrweH9\x00\x01"
            b"VM`ye`)G\x00\x01"
            b"\xff\xfe\xfd\xfc"
            b"Readable Title Slide\x00"
            b"Operating System Synchronization\x00"
        )

        upload_resp = self.client.post(
            '/api/files/upload',
            headers={"Authorization": f"Bearer {self.token}"},
            data={"file": (BytesIO(noisy_payload), "noisy.ppt")},
            content_type='multipart/form-data'
        )
        self.assertEqual(upload_resp.status_code, 201)
        file_id = upload_resp.get_json()["file"]["id"]

        preview_resp = self.client.get(
            f'/api/files/{file_id}/preview?token={self.token}'
        )
        self.assertEqual(preview_resp.status_code, 200)
        body = preview_resp.get_json()
        preview_html = body["preview"]["html"]
        self.assertNotIn("IDAT", preview_html)
        self.assertNotIn("p9dy9", preview_html)
        self.assertNotIn("QrweH9", preview_html)
        self.assertIn("Readable Title Slide", preview_html)

        insight_resp = self.client.get(
            f'/api/files/{file_id}/insights',
            headers={"Authorization": f"Bearer {self.token}"}
        )
        self.assertEqual(insight_resp.status_code, 200)
        insight_body = insight_resp.get_json()
        summary = insight_body["insights"]["summary"]
        self.assertNotIn("IDAT", summary)
        self.assertIn(insight_body["insights"]["source"], ["groq", "groq_unavailable"])
        if insight_body["insights"]["source"] == "groq_unavailable":
            self.assertIn("temporarily unavailable", summary)

    def test_mislabeled_ppt_with_pptx_payload_uses_real_presentation_preview(self):
        if not PPTX_AVAILABLE:
            self.skipTest("python-pptx is not available")

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "OS Process Synchronization"
        slide.placeholders[1].text = "Critical section and semaphore overview"

        stream = BytesIO()
        presentation.save(stream)
        stream.seek(0)

        upload_resp = self.client.post(
            '/api/files/upload',
            headers={"Authorization": f"Bearer {self.token}"},
            data={"file": (BytesIO(stream.getvalue()), "lecture.ppt")},
            content_type='multipart/form-data'
        )
        self.assertEqual(upload_resp.status_code, 201)
        file_id = upload_resp.get_json()["file"]["id"]

        preview_resp = self.client.get(
            f'/api/files/{file_id}/preview?token={self.token}'
        )
        self.assertEqual(preview_resp.status_code, 200)
        if 'application/json' in preview_resp.content_type:
            body = preview_resp.get_json()
            preview_html = body["preview"]["html"]
            self.assertIn("Slide 1", preview_html)
            self.assertIn("OS Process Synchronization", preview_html)
            self.assertIn("Critical section and semaphore overview", preview_html)
            self.assertNotIn("Legacy PowerPoint Preview", preview_html)
        else:
            self.assertIn('application/pdf', preview_resp.content_type)

    def test_preview_raw_mode_returns_native_stream_for_pptx(self):
        if not PPTX_AVAILABLE:
            self.skipTest("python-pptx is not available")

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Native PPTX Preview"
        slide.placeholders[1].text = "Open this in inline native mode"

        stream = BytesIO()
        presentation.save(stream)
        stream.seek(0)

        upload_resp = self.client.post(
            '/api/files/upload',
            headers={"Authorization": f"Bearer {self.token}"},
            data={"file": (BytesIO(stream.getvalue()), "native.pptx")},
            content_type='multipart/form-data'
        )
        self.assertEqual(upload_resp.status_code, 201)
        file_id = upload_resp.get_json()["file"]["id"]

        raw_resp = self.client.get(
            f'/api/files/{file_id}/preview?token={self.token}&mode=raw'
        )
        self.assertEqual(raw_resp.status_code, 200)
        self.assertIn('application/vnd.openxmlformats-officedocument.presentationml.presentation', raw_resp.content_type)


if __name__ == '__main__':
    unittest.main()