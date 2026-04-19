import json
import re
import zipfile
from collections import Counter
from io import BytesIO
from urllib.error import HTTPError, URLError
from urllib.request import Request, urlopen

from config import Config

try:
    from groq import Groq
    GROQ_SDK_AVAILABLE = True
except ImportError:
    GROQ_SDK_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    from pypdf import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False


class FileIntelligenceService:
    SUMMARY_MIN_WORDS = 50
    SUMMARY_MAX_WORDS = 100

    OOXML_NOISE_MARKERS = (
        '[Content_Types].xml',
        '_rels/.rels',
        'docProps/',
        'ppt/slideLayouts/',
        'ppt/slideMasters/',
        'ppt/theme/',
        'tableStyles.xml',
        'shapexml.xml',
        'downrev.xml',
    )

    TEXT_FILE_TYPES = {
        'txt', 'md', 'rst', 'rtf', 'json', 'xml', 'csv', 'tsv',
        'yaml', 'yml', 'toml', 'ini', 'cfg', 'conf',
        'py', 'js', 'ts', 'jsx', 'tsx', 'java', 'cpp', 'c', 'cs',
        'go', 'rs', 'rb', 'php', 'sh', 'bash', 'sql', 'html', 'css', 'vue', 'dat', 'hex',
    }

    STOP_WORDS = {
        'the', 'and', 'for', 'that', 'this', 'with', 'from', 'have', 'your', 'you', 'are',
        'was', 'were', 'will', 'would', 'there', 'their', 'about', 'into', 'after', 'before',
        'over', 'under', 'between', 'also', 'than', 'then', 'them', 'they', 'been', 'being',
        'file', 'files', 'data', 'info', 'information', 'content', 'document', 'project',
        'using', 'used', 'use', 'can', 'may', 'should', 'could', 'not', 'only', 'more',
        'most', 'very', 'some', 'any', 'one', 'two', 'three', 'new', 'old', 'such', 'when',
        'where', 'what', 'which', 'who', 'whom', 'why', 'how', 'has', 'had', 'its', 'our',
        'out', 'per', 'via', 'all', 'too', 'if', 'or', 'but', 'so', 'no',
    }

    SENSITIVE_KEYWORDS = {
        'credentials': ['password', 'passphrase', 'secret', 'token', 'api key', 'apikey', 'auth'],
        'personal': ['ssn', 'social security', 'address', 'phone', 'email', 'dob', 'date of birth'],
        'financial': ['credit card', 'card number', 'iban', 'bank', 'account number', 'payment'],
        'operations': ['roadmap', 'meeting', 'notes', 'confidential', 'internal', 'restricted'],
    }

    @staticmethod
    def _model_candidates():
        models = [Config.GROQ_MODEL, getattr(Config, 'GROQ_FALLBACK_MODEL', '')]
        return [m for i, m in enumerate(models) if m and m not in models[:i]]

    @staticmethod
    def build_insights(record, file_data):
        text, source = FileIntelligenceService._extract_text(record.file_type, file_data)
        text = (text or '').strip()

        local_insights = FileIntelligenceService._build_local_insights(record, text, source)

        if not Config.GROQ_API_KEY:
            return FileIntelligenceService._build_ai_unavailable_insights(
                record,
                text,
                source,
                "GROQ_API_KEY is missing.",
                local_insights,
            )

        groq_insights, groq_error = FileIntelligenceService._build_groq_insights(
            record, text, source, local_insights
        )
        if groq_insights:
            return groq_insights

        return FileIntelligenceService._build_ai_unavailable_insights(
            record,
            text,
            source,
            groq_error or "Groq request failed or returned invalid JSON.",
            local_insights,
        )

    @staticmethod
    def _build_ai_unavailable_insights(record, text, source, reason, local_insights):
        default_model = FileIntelligenceService._model_candidates()[0]
        return {
            'summary': 'AI insights are temporarily unavailable. Please verify Groq API configuration and try again.',
            'keywords': [],
            'tags': [record.file_type or 'file', 'ai-unavailable'],
            'sensitivity': 'unknown',
            'metrics': local_insights['metrics'],
            'source': 'groq_unavailable',
            'suggested_actions': [
                'Verify GROQ_API_KEY and GROQ_MODEL configuration',
                'Retry after confirming network access to Groq API',
            ],
            'engine': f'groq:{default_model}',
            'error': reason,
        }

    @staticmethod
    def _build_local_insights(record, text, source):
        word_count = len(re.findall(r"[A-Za-z0-9_']+", text))
        line_count = text.count('\n') + (1 if text else 0)
        character_count = len(text)

        keywords = FileIntelligenceService._extract_keywords(text)
        summary = FileIntelligenceService._summarize_text(text)
        tags = FileIntelligenceService._derive_tags(record.file_type, keywords, text)
        sensitivity = FileIntelligenceService._score_sensitivity(text, keywords)

        return {
            'summary': summary,
            'keywords': keywords,
            'tags': tags,
            'sensitivity': sensitivity,
            'metrics': {
                'words': word_count,
                'lines': line_count,
                'characters': character_count,
            },
            'source': source,
            'suggested_actions': FileIntelligenceService._suggest_actions(record.file_type, sensitivity, tags),
        }

    @staticmethod
    def _build_groq_insights(record, text, source, fallback_insights):
        prompt_text = text[:Config.GROQ_MAX_INPUT_CHARS].strip()
        if not prompt_text:
            return None, 'No readable text extracted for AI analysis.'

        payload = {
            'temperature': 0.1,
            'max_tokens': 260,
            'messages': [
                {
                    'role': 'system',
                    'content': (
                        'You are a security-focused file analysis assistant for a secure document manager. '
                        'Return ONLY strict JSON with keys: summary, keywords, tags, sensitivity, suggested_actions. '
                        '\n'
                        'SUMMARY REQUIREMENTS (critical):\n'
                        '- Must be EXACTLY 50-100 words, no more, no less (count carefully)\n'
                        '- For documents: describe topic, purpose, key insights, and security implications with specific details\n'
                        '- For images: describe visual content, type, composition, and any visible text/patterns with specifics\n'
                        '- For presentations: summarize main themes and key points discussed across slides\n'
                        '- For spreadsheets: describe data structure, subject matter, and purpose of analysis\n'
                        '- Avoid generic phrases like "appears to be" or "seems to"; use definitive language\n'
                        '- Include concrete observations, not filler text\n'
                        '- If OCR/image has visible text, mention key visible content\n'
                        '\n'
                        'KEYWORDS: lowercase, comma-separated, max 8, no generic terms\n'
                        'TAGS: lowercase, comma-separated, max 6\n'
                        'SENSITIVITY: exactly "low", "medium", or "high"\n'
                        'SUGGESTED_ACTIONS: 2-4 concise actionable items\n'
                        '\n'
                        'Do not include markdown, code fences, extra text, or explanations. Only JSON.'
                    ),
                },
                {
                    'role': 'user',
                    'content': (
                        f'Analyze this file for a security-focused dashboard. Be precise and specific in the summary (50-100 words exactly).\n\n'
                        f'Filename: {record.original_name}\n'
                        f'File type: {record.file_type}\n'
                        f'Preview source: {source}\n\n'
                        f'Content to analyze:\n{prompt_text}'
                    ),
                },
            ],
        }

        content = None
        used_model = None
        model_errors = []

        for model_name in FileIntelligenceService._model_candidates():
            request_body = {
                'model': model_name,
                'temperature': payload['temperature'],
                'max_tokens': payload['max_tokens'],
                'response_format': {'type': 'json_object'},
                'messages': payload['messages'],
            }

            request = Request(
                Config.GROQ_API_URL,
                data=json.dumps(request_body).encode('utf-8'),
                headers={
                    'Authorization': f'Bearer {Config.GROQ_API_KEY}',
                    'Content-Type': 'application/json',
                    'User-Agent': 'SecureFileManager/1.0',
                },
                method='POST',
            )

            try:
                with urlopen(request, timeout=20) as response:
                    response_payload = json.loads(response.read().decode('utf-8'))
                content = response_payload['choices'][0]['message']['content']
                used_model = model_name
                break
            except HTTPError as exc:
                try:
                    detail = exc.read().decode('utf-8', errors='ignore')
                except Exception:
                    detail = ''
                model_errors.append(f'HTTP {exc.code}: {detail[:600]}')

                # If rate limited on this model, try fallback model.
                if exc.code == 429 or 'rate limit' in detail.lower():
                    continue

                # Non-rate limit failure: do not keep retrying aggressively.
                break
            except (URLError, TimeoutError, ValueError, json.JSONDecodeError, OSError) as exc:
                model_errors.append(str(exc))
                break
            except (KeyError, IndexError, TypeError):
                model_errors.append('Groq API response did not contain chat content.')
                break

        if not content:
            return None, ' | '.join(model_errors) or 'Empty content returned by Groq API.'

        parsed = FileIntelligenceService._parse_ai_json(content)
        if not parsed:
            return None, 'Groq returned non-JSON output.'

        summary = str(parsed.get('summary') or fallback_insights['summary'])
        summary = FileIntelligenceService._normalize_document_summary(
            summary,
            prompt_text,
            fallback_insights.get('keywords', []),
            FileIntelligenceService.SUMMARY_MIN_WORDS,
            FileIntelligenceService.SUMMARY_MAX_WORDS,
        )
        keywords = FileIntelligenceService._normalize_string_list(parsed.get('keywords'), fallback_insights['keywords'])
        tags = FileIntelligenceService._normalize_string_list(parsed.get('tags'), fallback_insights['tags'])
        sensitivity = str(parsed.get('sensitivity') or fallback_insights['sensitivity']).lower().strip()
        if sensitivity not in {'low', 'medium', 'high'}:
            sensitivity = fallback_insights['sensitivity']
        suggested_actions = FileIntelligenceService._normalize_string_list(
            parsed.get('suggested_actions'),
            fallback_insights['suggested_actions']
        )

        return {
            'summary': summary,
            'keywords': keywords[:8],
            'tags': tags[:6],
            'sensitivity': sensitivity,
            'metrics': fallback_insights['metrics'],
            'source': 'groq',
            'suggested_actions': suggested_actions[:4],
            'engine': f'groq:{used_model or FileIntelligenceService._model_candidates()[0]}',
        }, None

    @staticmethod
    def _extract_text(file_type, file_data):
        normalized_type = (file_type or '').lower()

        try:
            if normalized_type in FileIntelligenceService.TEXT_FILE_TYPES:
                return file_data.decode('utf-8', errors='ignore'), 'text'

            if normalized_type in {'docx', 'doc'} and DOCX_AVAILABLE:
                return FileIntelligenceService._extract_docx_text(file_data), 'document'

            if normalized_type in {'xlsx', 'xls'} and OPENPYXL_AVAILABLE:
                return FileIntelligenceService._extract_excel_text(file_data), 'spreadsheet'

            if normalized_type in {'ppt', 'pptx'} and PPTX_AVAILABLE:
                if normalized_type == 'pptx' or FileIntelligenceService._looks_like_pptx_package(file_data):
                    try:
                        return FileIntelligenceService._extract_powerpoint_text(file_data), 'presentation'
                    except Exception:
                        pass

            if normalized_type == 'ppt':
                return FileIntelligenceService._extract_legacy_powerpoint_text(file_data), 'legacy_presentation'

            if normalized_type == 'pptx':
                return FileIntelligenceService._extract_ascii_runs(file_data), 'presentation_binary'

            if normalized_type in {'png', 'jpg', 'jpeg', 'gif', 'webp', 'bmp', 'tif', 'tiff', 'ico'} and PIL_AVAILABLE:
                return FileIntelligenceService._extract_image_text(file_data), 'image'

            if normalized_type == 'pdf':
                return FileIntelligenceService._extract_pdf_text(file_data), 'document'

        except Exception:
            pass

        return FileIntelligenceService._extract_ascii_runs(file_data), 'binary'

    @staticmethod
    def _extract_docx_text(file_data):
        doc = Document(BytesIO(file_data))
        parts = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text.strip())
        return '\n'.join(parts)

    @staticmethod
    def _extract_excel_text(file_data):
        workbook = load_workbook(BytesIO(file_data), data_only=True)
        parts = []
        for sheet_name in workbook.sheetnames[:5]:
            worksheet = workbook[sheet_name]
            parts.append(f'Sheet: {sheet_name}')
            for row_index, row in enumerate(worksheet.iter_rows(min_row=1, max_row=50, values_only=True)):
                values = [str(cell).strip() for cell in row[:10] if cell not in (None, '')]
                if values:
                    parts.append(' | '.join(values))
                if row_index >= 49:
                    break
        return '\n'.join(parts)

    @staticmethod
    def _extract_powerpoint_text(file_data):
        presentation = Presentation(BytesIO(file_data))
        parts = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if slide_parts:
                parts.append(f'Slide {slide_index}: ' + ' '.join(slide_parts))
        return '\n'.join(parts)

    @staticmethod
    def _extract_legacy_powerpoint_text(file_data):
        """Best-effort text extraction for binary .ppt files."""
        runs = FileIntelligenceService._extract_meaningful_text_runs(file_data)
        if not runs:
            return 'Legacy PowerPoint file detected. No readable text was extracted.'

        return '\n'.join(runs)

    @staticmethod
    def _extract_image_text(file_data):
        try:
            image = Image.open(BytesIO(file_data))
            width, height = image.size
            mode = image.mode
            
            # Analyze image properties for better description
            format_name = image.format or 'Unknown'
            
            # Get dominant colors
            colors = []
            try:
                converted = image.convert('RGB')
                pixels = list(converted.getdata())
                if pixels:
                    from collections import Counter
                    color_counter = Counter(pixels)
                    dominant = color_counter.most_common(1)[0][0]
                    colors = f'dominant color RGB{dominant}'
            except:
                pass
            
            # Estimate image type by analyzing content
            image_type = FileIntelligenceService._classify_image(image)
            
            # Build comprehensive description
            description = f'Image: {width}x{height} pixels, {format_name} format, {mode} color mode. Type: {image_type}.'
            if colors:
                description += f' {colors}.'
            
            # Add file size info
            size_kb = len(file_data) / 1024
            description += f' File size: {size_kb:.1f} KB.'
            
            return description
        except Exception:
            return 'Image file detected but unable to extract properties for analysis.'

    @staticmethod
    def _extract_pdf_text(file_data):
        if PDF_AVAILABLE:
            try:
                reader = PdfReader(BytesIO(file_data))
                parts = []
                for page in reader.pages[:12]:
                    extracted = page.extract_text() or ''
                    extracted = re.sub(r'\s+', ' ', extracted).strip()
                    if extracted:
                        parts.append(extracted)
                if parts:
                    return '\n'.join(parts)
            except Exception:
                pass

        runs = FileIntelligenceService._extract_meaningful_text_runs(file_data)
        if runs:
            return '\n'.join(runs)

        return 'PDF file detected. No readable text was extracted for AI analysis.'

    @staticmethod
    def _extract_ascii_runs(file_data):
        if not isinstance(file_data, (bytes, bytearray)):
            return ''
        runs = re.findall(rb'[\x20-\x7E]{4,}', file_data)
        if not runs:
            return ''
        text = ' '.join(run.decode('utf-8', errors='ignore') for run in runs[:200])
        return re.sub(r'\s+', ' ', text)

    @staticmethod
    def _classify_image(image):
        """Classify image type based on properties."""
        try:
            width, height = image.size
            aspect_ratio = width / max(height, 1)
            
            # Check if image is mostly monochrome or has color
            if image.mode in ('L', 'LA', '1'):
                return 'grayscale/monochrome'
            
            # Estimate based on aspect ratio
            if aspect_ratio > 2.5:
                return 'wide format (banner/header)'
            elif aspect_ratio < 0.4:
                return 'portrait/vertical'
            elif abs(aspect_ratio - 1.0) < 0.1:
                return 'square (icon/thumbnail)'
            else:
                return 'landscape/standard'
        except:
            return 'general image'

    @staticmethod
    def _extract_meaningful_text_runs(file_data):
        if not isinstance(file_data, (bytes, bytearray)):
            return []

        raw_runs = re.findall(rb'[\x20-\x7E]{12,}', file_data)
        cleaned_runs = []

        for raw_run in raw_runs:
            text = raw_run.decode('utf-8', errors='ignore').strip()
            text = re.sub(r'\s+', ' ', text)
            if not text:
                continue

            if len(text) > 140:
                text = text[:140].rstrip()

            tokens = re.findall(r"[A-Za-z][A-Za-z0-9'\-]{2,}", text)
            if len(tokens) < 2:
                continue

            alpha_count = sum(1 for ch in text if ch.isalpha())
            printable_count = sum(1 for ch in text if ch.isprintable())
            if printable_count == 0:
                continue

            alpha_ratio = alpha_count / max(len(text), 1)
            if alpha_ratio < 0.55:
                continue

            # Keep mostly human-readable lines, not binary fragments.
            if any(fragment in text for fragment in ('IDAT', 'IHDR', 'pHYs', 'gAMA', 'sRGB', 'tEXt')):
                continue

            if any(marker in text for marker in FileIntelligenceService.OOXML_NOISE_MARKERS):
                continue

            if re.search(r'\b(?:ppt|_rels|docProps|word|xl)/[^\s]{2,}\b', text):
                continue

            if text.endswith('PK'):
                continue

            cleaned_runs.append(text)

        # De-duplicate while preserving order.
        unique_runs = []
        for run in cleaned_runs:
            if run not in unique_runs:
                unique_runs.append(run)

        return unique_runs[:24]

    @staticmethod
    def _looks_like_pptx_package(file_data):
        if not isinstance(file_data, (bytes, bytearray)):
            return False

        if not file_data.startswith(b'PK'):
            return False

        try:
            with zipfile.ZipFile(BytesIO(file_data)) as archive:
                names = set(archive.namelist())
                if '[Content_Types].xml' not in names:
                    return False

                has_ppt_payload = any(name.startswith('ppt/') for name in names)
                has_relationships = '_rels/.rels' in names
                return has_ppt_payload and has_relationships
        except (OSError, zipfile.BadZipFile, RuntimeError):
            return False

    @staticmethod
    def _extract_keywords(text, limit=8):
        words = re.findall(r"[A-Za-z0-9_']+", (text or '').lower())
        filtered = [word for word in words if len(word) > 2 and word not in FileIntelligenceService.STOP_WORDS]
        if not filtered:
            return []
        counts = Counter(filtered)
        return [word for word, _ in counts.most_common(limit)]

    @staticmethod
    def _summarize_text(text):
        cleaned = re.sub(r'\s+', ' ', (text or '').strip())
        if not cleaned:
            return 'No readable text could be extracted from this file.'

        summary = FileIntelligenceService._build_document_oriented_summary(cleaned, [])
        return FileIntelligenceService._normalize_document_summary(
            summary,
            cleaned,
            [],
            FileIntelligenceService.SUMMARY_MIN_WORDS,
            FileIntelligenceService.SUMMARY_MAX_WORDS,
        )

    @staticmethod
    def _word_count(text):
        return len(re.findall(r"[A-Za-z0-9_']+", text or ''))

    @staticmethod
    def _trim_to_word_limit(text, max_words):
        words = re.findall(r"\S+", text or '')
        if len(words) <= max_words:
            return (text or '').strip()

        trimmed = ' '.join(words[:max_words]).strip()
        return trimmed.rstrip(',.') + '.'

    @staticmethod
    def _build_document_oriented_summary(text, keywords):
        cleaned = re.sub(r'\s+', ' ', (text or '').strip())
        if not cleaned:
            return ''

        sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', cleaned) if s.strip()]
        collected = []
        for sentence in sentences:
            candidate = sentence
            if len(candidate) < 20:
                continue
            if candidate not in collected:
                collected.append(candidate)
            combined = ' '.join(collected)
            if FileIntelligenceService._word_count(combined) >= 70:
                break

        if not collected:
            snippet = cleaned[:700].strip()
            collected.append(snippet)

        summary = ' '.join(collected)

        lowered = summary.lower()
        if 'document' not in lowered and 'presentation' not in lowered and 'file' not in lowered:
            summary = f'This document discusses {summary[0].lower() + summary[1:] if len(summary) > 1 else summary} '

        if keywords:
            top = ', '.join(keywords[:3])
            hint = f' Key topics include {top}. '
            if hint.strip().lower() not in summary.lower():
                summary = (summary + hint).strip()

        return summary

    @staticmethod
    def _normalize_document_summary(summary, source_text, keywords, min_words, max_words):
        text = re.sub(r'\s+', ' ', (summary or '').strip())
        if not text:
            text = FileIntelligenceService._build_document_oriented_summary(source_text, keywords)

        word_count = FileIntelligenceService._word_count(text)
        
        # If too long, trim to max
        if word_count > max_words:
            text = FileIntelligenceService._trim_to_word_limit(text, max_words)
            word_count = FileIntelligenceService._word_count(text)

        # If too short, supplement with meaningful content
        if word_count < min_words:
            supplement = FileIntelligenceService._build_document_oriented_summary(source_text, keywords)
            if supplement and supplement.lower() not in text.lower():
                text = (text + ' ' + supplement).strip()
                word_count = FileIntelligenceService._word_count(text)

            # If still short, add security/practical guidance
            if word_count < min_words:
                additions = [
                    'Ensure encrypted storage and access controls.',
                    'Review sensitivity level before sharing.',
                    'Maintain audit logs for access tracking.',
                    'Consider data retention policies.',
                    'Verify recipient permissions align with content classification.',
                ]
                for addition in additions:
                    candidate = (text + ' ' + addition).strip()
                    if FileIntelligenceService._word_count(candidate) <= max_words:
                        text = candidate
                        word_count = FileIntelligenceService._word_count(text)
                        if word_count >= min_words:
                            break

            # Final trim to respect max_words
            if word_count > max_words:
                text = FileIntelligenceService._trim_to_word_limit(text, max_words)

        return text.strip()

    @staticmethod
    def _derive_tags(file_type, keywords, text):
        tags = []
        normalized_type = (file_type or 'unknown').lower()

        type_tags = {
            'pdf': 'document',
            'doc': 'document',
            'docx': 'document',
            'txt': 'text',
            'md': 'text',
            'csv': 'spreadsheet',
            'tsv': 'spreadsheet',
            'xls': 'spreadsheet',
            'xlsx': 'spreadsheet',
            'ppt': 'presentation',
            'pptx': 'presentation',
            'png': 'image',
            'jpg': 'image',
            'jpeg': 'image',
            'gif': 'image',
            'webp': 'image',
            'mp4': 'video',
            'mp3': 'audio',
            'zip': 'archive',
            'json': 'data',
            'xml': 'data',
            'yaml': 'data',
            'yml': 'data',
            'toml': 'data',
        }
        if normalized_type in type_tags:
            tags.append(type_tags[normalized_type])

        if normalized_type == 'ppt':
            tags.append('legacy-ppt')
        elif normalized_type == 'pptx':
            tags.append('presentation-preview')

        for keyword in keywords[:4]:
            if keyword not in tags:
                tags.append(keyword)

        if re.search(r'confidential|restricted|internal', text or '', re.IGNORECASE):
            tags.append('sensitive')
        if re.search(r'project|plan|roadmap|milestone', text or '', re.IGNORECASE):
            tags.append('planning')
        if re.search(r'invoice|payment|budget|finance', text or '', re.IGNORECASE):
            tags.append('finance')

        unique_tags = []
        for tag in tags:
            if tag not in unique_tags:
                unique_tags.append(tag)
        return unique_tags[:6]

    @staticmethod
    def _score_sensitivity(text, keywords):
        haystack = f"{text or ''} {' '.join(keywords or [])}"
        score = 0
        for patterns in FileIntelligenceService.SENSITIVE_KEYWORDS.values():
            for pattern in patterns:
                if re.search(pattern, haystack, re.IGNORECASE):
                    score += 1
                    break

        if score >= 4:
            return 'high'
        if score >= 2:
            return 'medium'
        return 'low'

    @staticmethod
    def _suggest_actions(file_type, sensitivity, tags):
        actions = ['Keep encrypted at rest', 'Share only with least-privilege access']
        if sensitivity == 'high':
            actions.insert(0, 'Review before sharing externally')
        if 'finance' in tags or 'sensitive' in tags:
            actions.append('Consider expiring access links or read-only sharing')
        if file_type in {'csv', 'xlsx', 'json', 'xml'}:
            actions.append('Check for structured data leaks before syncing to cloud storage')
        return actions[:4]

    @staticmethod
    def _parse_ai_json(content):
        if not content:
            return None

        text = str(content).strip()
        if text.startswith('```'):
            text = re.sub(r'^```(?:json)?\s*', '', text, flags=re.IGNORECASE)
            text = re.sub(r'\s*```$', '', text)

        start = text.find('{')
        end = text.rfind('}')
        if start == -1 or end == -1 or end <= start:
            return None

        try:
            return json.loads(text[start:end + 1])
        except json.JSONDecodeError:
            return None

    @staticmethod
    def _normalize_string_list(value, fallback):
        if isinstance(value, list):
            items = [str(item).strip() for item in value if str(item).strip()]
            if items:
                return items
        return list(fallback)