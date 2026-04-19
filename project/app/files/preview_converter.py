import base64
import json
import csv
import os
import re
import tempfile
import zipfile
from io import StringIO, BytesIO
from pathlib import Path

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    import pythoncom
    import win32com.client
    WIN32COM_AVAILABLE = True
except ImportError:
    WIN32COM_AVAILABLE = False


class PreviewConverter:
    OOXML_NOISE_MARKERS = (
        '[Content_Types].xml',
        '_rels/.rels',
        'docProps/',
        'ppt/slideLayouts/',
        'ppt/slideMasters/',
        'ppt/theme/',
        'tableStyles.xml',
        'shapexml.xml',
        'downrev.xml',
    )

    @staticmethod
    def convert_file(file_data, filename, file_type):
        """Convert file to preview HTML based on file type"""
        normalized_type = (file_type or '').lower()

        text_preview_types = {
            'txt', 'md', 'json', 'xml', 'csv', 'yaml', 'yml', 'html', 'css', 'js', 'py', 'java', 'cpp', 'c', 'sql',
            'log', 'rst', 'rtf', 'tsv', 'toml', 'ini', 'cfg', 'conf',
            'ts', 'jsx', 'tsx', 'vue', 'go', 'rs', 'rb', 'php', 'sh', 'bash', 'cs', 'dat', 'hex'
        }
        
        # Text-based files - always convertible
        if normalized_type in text_preview_types:
            return PreviewConverter._preview_text(file_data, filename, normalized_type)
        
        # Office documents - with library checks
        if normalized_type in ['docx', 'doc'] and DOCX_AVAILABLE:
            result = PreviewConverter._preview_docx(file_data, filename)
            if result:
                return result
                
        if normalized_type in ['xlsx', 'xls'] and OPENPYXL_AVAILABLE:
            result = PreviewConverter._preview_excel(file_data, filename)
            if result:
                return result
                
        # Some files are mislabeled as .ppt but actually contain a PPTX/OOXML package.
        if normalized_type in ['ppt', 'pptx'] and PPTX_AVAILABLE:
            if normalized_type == 'pptx' or PreviewConverter._looks_like_pptx_package(file_data):
                result = PreviewConverter._preview_powerpoint(file_data, filename)
                if result:
                    return result

        # Legacy .ppt fallback when OOXML parsing isn't possible.
        if normalized_type == 'ppt':
            result = PreviewConverter._preview_legacy_powerpoint(file_data, filename)
            if result:
                return result

        # Fallback for unknown or partially supported presentation payloads.
        if normalized_type in ['ppt', 'pptx']:
            result = PreviewConverter._preview_legacy_powerpoint(file_data, filename)
            if result:
                return result
        
        # Return None to trigger fallback in route
        return None

    @staticmethod
    def convert_presentation_to_pdf_with_ms_office(file_data, filename, file_type):
        """Convert PPT/PPTX to PDF using installed Microsoft PowerPoint on Windows."""
        normalized_type = (file_type or '').lower()
        if normalized_type not in {'ppt', 'pptx'}:
            return None

        if os.name != 'nt' or not WIN32COM_AVAILABLE:
            return None

        effective_type = normalized_type
        if normalized_type == 'ppt' and PreviewConverter._looks_like_pptx_package(file_data):
            effective_type = 'pptx'

        safe_stem = Path(filename or 'presentation').stem or 'presentation'
        safe_stem = re.sub(r'[^A-Za-z0-9._-]+', '_', safe_stem)

        presentation = None
        app = None

        with tempfile.TemporaryDirectory(prefix='sfm_ms_preview_') as tmp_dir:
            input_path = os.path.join(tmp_dir, f'{safe_stem}.{effective_type}')
            output_path = os.path.join(tmp_dir, f'{safe_stem}.pdf')

            with open(input_path, 'wb') as handle:
                handle.write(file_data)

            try:
                pythoncom.CoInitialize()
                app = win32com.client.DispatchEx('PowerPoint.Application')

                # Open read-only, hidden window.
                presentation = app.Presentations.Open(input_path, ReadOnly=1, WithWindow=0)

                # 32 => ppSaveAsPDF
                presentation.SaveAs(output_path, 32)

                if not os.path.exists(output_path):
                    return None

                with open(output_path, 'rb') as pdf_handle:
                    return pdf_handle.read()
            except Exception:
                return None
            finally:
                try:
                    if presentation is not None:
                        presentation.Close()
                except Exception:
                    pass

                try:
                    if app is not None:
                        app.Quit()
                except Exception:
                    pass

                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass

    @staticmethod
    def _preview_text(file_data, filename, file_type):
        """Preview text-based files"""
        try:
            text = file_data.decode('utf-8', errors='ignore')
            
            # Handle JSON pretty printing
            if file_type == 'json':
                try:
                    parsed = json.loads(text)
                    text = json.dumps(parsed, indent=2)
                except:
                    pass
            
            # Truncate very long files
            if len(text) > 50000:
                text = text[:50000] + '\n\n... [File truncated for preview] ...'
            
            # Escape HTML
            safe_text = _escape_html(text)
            
            return {
                'type': 'text',
                'html': f'<pre style="white-space:pre-wrap;max-height:70vh;overflow:auto;background:#0f1320;padding:14px;border-radius:8px;font-size:.85rem;line-height:1.5;font-family:monospace;">{safe_text}</pre>',
                'filename': filename
            }
        except Exception as e:
            return None

    @staticmethod
    def _preview_docx(file_data, filename):
        """Preview Word documents"""
        if not DOCX_AVAILABLE:
            return None
        
        try:
            doc = Document(BytesIO(file_data))
            
            html_parts = ['<div style="padding:20px;max-height:70vh;overflow:auto;background:#f9f9f9;">']
            
            para_count = 0
            for para in doc.paragraphs:
                if para.text.strip():
                    style = para.style.name if para.style else 'Normal'
                    html_parts.append(f'<p style="margin:10px 0;line-height:1.6;">{_escape_html(para.text)}</p>')
                    para_count += 1
                    if para_count > 100:
                        html_parts.append('<p style="color:#999;font-style:italic;">... [Content truncated] ...</p>')
                        break
            
            for table in doc.tables:
                html_parts.append('<table style="width:100%;border-collapse:collapse;margin:20px 0;border:1px solid #ddd;">')
                for row_idx, row in enumerate(table.rows):
                    html_parts.append('<tr>')
                    for cell in row.cells:
                        bg = '#e8e8e8' if row_idx == 0 else 'white'
                        html_parts.append(f'<td style="border:1px solid #ddd;padding:8px;background:{bg};">{_escape_html(cell.text)}</td>')
                    html_parts.append('</tr>')
                html_parts.append('</table>')
            
            html_parts.append('</div>')
            
            return {
                'type': 'document',
                'html': ''.join(html_parts),
                'filename': filename
            }
        except Exception as e:
            print(f"DOCX Error: {str(e)}")
            return None

    @staticmethod
    def _preview_excel(file_data, filename):
        """Preview Excel spreadsheets"""
        if not OPENPYXL_AVAILABLE:
            return None
        
        try:
            wb = load_workbook(BytesIO(file_data), data_only=True)
            
            html_parts = ['<div style="padding:20px;max-height:70vh;overflow:auto;background:#f9f9f9;"><div style="overflow-x:auto;">']
            
            for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                if sheet_idx >= 5:  # Limit to first 5 sheets
                    break
                    
                ws = wb[sheet_name]
                html_parts.append(f'<h3 style="margin-top:20px;margin-bottom:10px;font-size:1rem;">{_escape_html(sheet_name)}</h3>')
                html_parts.append('<table style="border-collapse:collapse;margin:10px 0;border:1px solid #ddd;font-size:0.9rem;">')
                
                row_count = 0
                for row_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=100, values_only=True)):
                    if row_count > 50:
                        html_parts.append('<tr><td colspan="10" style="padding:10px;text-align:center;color:#999;">... [More rows] ...</td></tr>')
                        break
                    
                    html_parts.append('<tr>')
                    for cell_idx, cell_val in enumerate(row[:10]):  # Limit to 10 columns
                        val = str(cell_val) if cell_val is not None else ''
                        bg = '#e8e8e8' if row_idx == 0 else 'white'
                        html_parts.append(f'<td style="border:1px solid #ddd;padding:6px;background:{bg};min-width:60px;">{_escape_html(val)}</td>')
                    html_parts.append('</tr>')
                    row_count += 1
                
                html_parts.append('</table>')
            
            html_parts.append('</div></div>')
            
            return {
                'type': 'spreadsheet',
                'html': ''.join(html_parts),
                'filename': filename
            }
        except Exception as e:
            print(f"Excel Error: {str(e)}")
            return None

    @staticmethod
    def _preview_powerpoint(file_data, filename):
        """Preview PowerPoint presentations"""
        if not PPTX_AVAILABLE:
            return None
        
        try:
            prs = Presentation(BytesIO(file_data))
            
            html_parts = ['<div style="padding:20px;max-height:70vh;overflow:auto;background:#f9f9f9;">']
            
            for slide_idx, slide in enumerate(prs.slides):
                if slide_idx >= 15:  # Limit to first 15 slides
                    html_parts.append(f'<p style="text-align:center;color:#999;margin:20px 0;">... [{len(prs.slides) - 15} more slides] ...</p>')
                    break
                    
                html_parts.append(f'<div style="margin-bottom:20px;border:1px solid #ddd;padding:15px;border-radius:8px;background:white;">')
                html_parts.append(f'<h4 style="margin-top:0;color:#333;font-size:0.95rem;">Slide {slide_idx + 1}</h4>')
                
                # Extract text from shapes
                slide_text = []
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'text') and shape.text.strip():
                            slide_text.append(shape.text)
                    except:
                        pass
                
                if slide_text:
                    for text in slide_text:
                        html_parts.append(f'<p style="margin:8px 0;line-height:1.5;">{_escape_html(text)}</p>')
                else:
                    html_parts.append('<p style="color:#999;margin:8px 0;font-style:italic;">[Slide with images or minimal text]</p>')
                
                html_parts.append('</div>')
            
            html_parts.append('</div>')
            
            return {
                'type': 'presentation',
                'html': ''.join(html_parts),
                'filename': filename
            }
        except Exception as e:
            print(f"PowerPoint Error: {str(e)}")
            return None

    @staticmethod
    def _preview_legacy_powerpoint(file_data, filename):
        """Best-effort preview for legacy .ppt by extracting printable text runs."""
        try:
            if not isinstance(file_data, (bytes, bytearray)):
                return None

            runs = PreviewConverter._extract_meaningful_text_runs(file_data)
            if not runs:
                return {
                    'type': 'presentation',
                    'html': (
                        '<div style="padding:20px;max-height:70vh;overflow:auto;background:#f9f9f9;">'
                        '<p style="margin-bottom:10px;">Legacy PowerPoint (.ppt) detected.</p>'
                        '<p style="color:#666;">No extractable text was found in this file. '
                        'Download to open in PowerPoint for full fidelity.</p>'
                        '</div>'
                    ),
                    'filename': filename
                }

            text_lines = runs

            html_items = ''.join(
                f'<li style="margin:6px 0;line-height:1.45;">{_escape_html(line)}</li>'
                for line in text_lines[:80]
            )

            html = (
                '<div style="padding:20px;max-height:70vh;overflow:auto;background:#f9f9f9;">'
                '<p style="margin-bottom:12px;"><strong>Legacy PowerPoint Preview (.ppt)</strong></p>'
                '<p style="margin-bottom:14px;color:#666;">'
                'This is a best-effort text extraction preview for an older .ppt file. '
                'Layout, images, and animations are not rendered here.'
                '</p>'
                '<ul style="padding-left:20px;">'
                f'{html_items}'
                '</ul>'
                '</div>'
            )

            return {
                'type': 'presentation',
                'html': html,
                'filename': filename
            }
        except Exception as e:
            print(f"Legacy PowerPoint Error: {str(e)}")
            return None

    @staticmethod
    def _extract_meaningful_text_runs(file_data):
        if not isinstance(file_data, (bytes, bytearray)):
            return []

        raw_runs = re.findall(rb'[\x20-\x7E]{12,}', file_data)
        cleaned_runs = []

        for raw_run in raw_runs:
            text = raw_run.decode('utf-8', errors='ignore').strip()
            text = re.sub(r'\s+', ' ', text)
            if not text:
                continue

            if len(text) > 140:
                text = text[:140].rstrip()

            tokens = re.findall(r"[A-Za-z][A-Za-z0-9'\-]{2,}", text)
            if len(tokens) < 2:
                continue

            alpha_count = sum(1 for ch in text if ch.isalpha())
            alpha_ratio = alpha_count / max(len(text), 1)
            if alpha_ratio < 0.55:
                continue

            if any(fragment in text for fragment in ('IDAT', 'IHDR', 'pHYs', 'gAMA', 'sRGB', 'tEXt')):
                continue

            # Filter common OOXML package internals that are not user-facing slide text.
            if any(marker in text for marker in PreviewConverter.OOXML_NOISE_MARKERS):
                continue

            # Skip archive-like path fragments such as "ppt/slides/slide1.xmlPK".
            if re.search(r'\b(?:ppt|_rels|docProps|word|xl)/[^\s]{2,}\b', text):
                continue

            if text.endswith('PK'):
                continue

            cleaned_runs.append(text)

        unique_runs = []
        for run in cleaned_runs:
            if run not in unique_runs:
                unique_runs.append(run)

        return unique_runs[:24]

    @staticmethod
    def _looks_like_pptx_package(file_data):
        if not isinstance(file_data, (bytes, bytearray)):
            return False

        # OOXML containers are ZIP payloads.
        if not file_data.startswith(b'PK'):
            return False

        try:
            with zipfile.ZipFile(BytesIO(file_data)) as archive:
                names = set(archive.namelist())
                if '[Content_Types].xml' not in names:
                    return False

                # PPTX packages include ppt/ directories and relationships metadata.
                has_ppt_payload = any(name.startswith('ppt/') for name in names)
                has_relationships = '_rels/.rels' in names
                return has_ppt_payload and has_relationships
        except (OSError, zipfile.BadZipFile, RuntimeError):
            return False


def _escape_html(text):
    """Escape HTML special characters"""
    if not text:
        return ''
    return (str(text).replace('&', '&amp;')
                     .replace('<', '&lt;')
                     .replace('>', '&gt;')
                     .replace('"', '&quot;')
                     .replace("'", '&#39;'))
